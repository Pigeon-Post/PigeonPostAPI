# database/lecture_content.py
import io
from openai import OpenAI
from datetime import datetime
from config.config import Config
from database.google_drive import GoogleDriveUploader
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import json

class LectureContentGenerator:
    def __init__(self):
        self.client = OpenAI(api_key=Config.OPENAI_API_KEY)
        self.drive_uploader = GoogleDriveUploader()
        self.MAX_TTS_LENGTH = 3000

    def generate_podcast_script(self, lecture_content: str, course_title: str, lecture_title: str) -> str:
        """Generate a single, flowing podcast script"""
        prompt = f"""
        Create a concise lecture (maximum 1500 words) for {course_title}, {lecture_title}.
        Include:
        1. Brief welcome (1 sentence)
        2. Clear explanation of 3-4 key concepts
        3. Brief conclusion (1 sentence)

        Content to cover:
        {lecture_content}
        """

        response = self.client.chat.completions.create(
            model=Config.LLM_MODEL,
            messages=[
                {"role": "system", "content": "You are an engaging professor. Keep explanations brief and focused."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=1500
        )

        return response.choices[0].message.content

    def extract_lecture_structure(self, lecture_content: str) -> dict:
        """Extract structured content for slides"""
        system_prompt = "You are a teaching assistant that organizes lecture content into clear, structured slides."
        
        user_prompt = f"""Extract the main points from this lecture and format them as slides.
        Return your response in this exact JSON structure:
        {{
            "title": "Lecture title",
            "sections": [
                {{
                    "title": "Section title",
                    "key_points": ["point 1", "point 2", "point 3"],
                    "details": "Detailed explanation"
                }}
            ]
        }}

        Lecture content:
        {lecture_content}
        """

        response = self.client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.5,
            response_format={"type": "json_object"}
        )
        
        return json.loads(response.choices[0].message.content)

    def create_powerpoint(self, lecture_data: dict) -> bytes:
        """Create PowerPoint slides and return as bytes"""
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = lecture_data["title"]
        subtitle.text = f"Generated on: {datetime.now().strftime('%Y-%m-%d')}"

        # Content slides
        for section in lecture_data["sections"]:
            # Section title and key points
            content_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_slide_layout)
            
            title = slide.shapes.title
            title.text = section["title"]
            
            # Add key points
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.text = "Key Points:"
            
            for point in section["key_points"]:
                p = tf.add_paragraph()
                p.text = f"• {point}"
                p.level = 1
            
            # Details slide
            detail_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(detail_slide_layout)
            
            title = slide.shapes.title
            title.text = f"{section['title']} - Details"
            
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.text = section["details"]

        output_buffer = io.BytesIO()
        prs.save(output_buffer)
        output_buffer.seek(0)
        return output_buffer.getvalue()

    def generate_content(self, lecture_content: str, course_title: str, lecture_title: str, email: str = None) -> dict:
        """Generate both podcast and slides"""
        try:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            safe_prefix = f"{course_title}_{lecture_title}".lower()
            safe_prefix = "".join(c for c in safe_prefix if c.isalnum() or c in (' ', '_')).replace(' ', '_')

            # Create folder
            folder_name = f"{safe_prefix}_{timestamp}"
            folder_id = self.drive_uploader.create_folder(folder_name)

            results = {
                "status": "success",
                "audio": None,
                "slides": None
            }

            # Generate and upload podcast
            print("Generating podcast script...")
            script = self.generate_podcast_script(lecture_content, course_title, lecture_title)
            
            print("Generating audio...")
            audio_response = self.client.audio.speech.create(
                model="tts-1-hd",
                voice="nova",
                input=script[:4000],
                speed=0.97
            )

            print("Uploading podcast...")
            audio_stream = io.BytesIO(audio_response.content)
            audio_filename = f"{safe_prefix}_audio_{timestamp}.mp3"
            results["audio"] = self.drive_uploader.upload_stream(
                audio_stream,
                audio_filename,
                'audio/mpeg',
                folder_id,
                email
            )

            # Generate and upload PowerPoint
            print("Generating PowerPoint slides...")
            lecture_data = self.extract_lecture_structure(lecture_content)
            slides_pptx = self.create_powerpoint(lecture_data)
            
            print("Uploading slides...")
            slides_stream = io.BytesIO(slides_pptx)
            slides_filename = f"{safe_prefix}_slides_{timestamp}.pptx"
            results["slides"] = self.drive_uploader.upload_stream(
                slides_stream,
                slides_filename,
                'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                folder_id,
                email
            )

            # Check for success and include URLs in response
            if results["audio"]["status"] == "success" and results["slides"]["status"] == "success":
                return {
                    "status": "success",
                    "audio_link": results["audio"]["web_link"],
                    "slides_link": results["slides"]["web_link"],
                    "message": "Podcast and PowerPoint slides generated successfully!"
                }
            else:
                return {
                    "status": "error",
                    "message": "Some uploads failed",
                    "audio_link": results["audio"].get("web_link"),
                    "slides_link": results["slides"].get("web_link")
                }

        except Exception as e:
            return {
                "status": "error",
                "message": f"Failed to generate content: {str(e)}"
            }